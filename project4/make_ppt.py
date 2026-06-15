"""
Generate Defense Presentation PPT for Project 4: Air-Conditioning Power Load Forecasting
==========================================================================================
Uses the CQU template backgrounds while creating entirely new content focused on
methodology and experimental findings (做了什么).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
import datetime

# Paths
BASE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = "/home/hhl123/projects/答辩2.pptx"
OUTPUT = os.path.join(BASE, "output", "空调负荷预测_答辩汇报.pptx")
FIGURES = os.path.join(BASE, "output")

# ======================== Color Palette ========================
CQU_BLUE   = RGBColor(0x1A, 0x3C, 0x6E)   # 重大蓝
CQU_RED    = RGBColor(0xC4, 0x1A, 0x23)    # 重大红
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE  = RGBColor(0x06, 0x5A, 0x82)
ACCENT_GREEN = RGBColor(0x2C, 0x5F, 0x2D)
ACCENT_RED   = RGBColor(0xF9, 0x61, 0x67)
TITLE_BG   = RGBColor(0x1A, 0x3C, 0x6E)

# ======================== Helpers ========================
# Module-level slide dimensions (set in build_ppt)
SW = None
SH = None

def slide_w(emu_val):
    """Convert cm to EMU based on actual slide width."""
    return Emu(emu_val)

def set_font(run, name="微软雅黑", size=Pt(14), bold=False, color=DARK_GRAY, east_asian=None):
    """Set font properties for a run."""
    run.font.size = size
    run.font.name = name
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    # Set East Asian font
    from pptx.oxml.ns import qn
    ea = east_asian or name
    rPr = run._r.get_or_add_rPr()
    # Remove existing ea elements
    for child in list(rPr):
        if child.tag == qn('a:ea'):
            rPr.remove(child)
    ea_elem = rPr.makeelement(qn('a:ea'), {})
    ea_elem.set('typeface', ea)
    rPr.append(ea_elem)

def add_textbox(slide, left, top, width, height, text, font_size=Pt(14),
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                font_name="微软雅黑", line_spacing=1.2):
    """Add a simple text box to a slide."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(2)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size.pt * line_spacing)
    run = p.add_run()
    run.text = text
    set_font(run, name=font_name, size=font_size, bold=bold, color=color)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=Pt(13),
                          color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="微软雅黑",
                          line_spacing=1.4):
    """Add a text box with multiple paragraphs (one per line), supporting bold markers."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, bold, clr, fsize = line_info, False, color, font_size
        else:
            text = line_info[0]
            bold = line_info[1] if len(line_info) > 1 else False
            clr = line_info[2] if len(line_info) > 2 else color
            fsize = line_info[3] if len(line_info) > 3 else font_size

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(3)
        p.line_spacing = Pt(fsize.pt * line_spacing)
        run = p.add_run()
        run.text = text
        set_font(run, name=font_name, size=fsize, bold=bold, color=clr)
    return txBox

def add_slide_title(slide, title_text, subtitle_text=None):
    """Add a standardized title bar at the top of a content slide."""
    # Title background bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        SW, Emu(Cm(1.8))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_BG
    shape.line.fill.background()

    # Title text
    add_textbox(slide, Cm(1.5), Cm(0.25), Cm(22), Cm(1.3),
                title_text, font_size=Pt(26), bold=True, color=WHITE,
                alignment=PP_ALIGN.LEFT, font_name="微软雅黑")

    if subtitle_text:
        add_textbox(slide, Cm(1.5), Cm(1.05), Cm(22), Cm(0.6),
                    subtitle_text, font_size=Pt(12), bold=False,
                    color=RGBColor(0xCC, 0xCC, 0xCC),
                    alignment=PP_ALIGN.LEFT, font_name="微软雅黑")

    # Bottom accent line
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(Cm(1.8)),
        SW, Emu(Cm(0.06))
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = CQU_RED
    shape2.line.fill.background()

def add_image_safe(slide, path, left, top, width, height=None, caption=None):
    """Add an image if it exists, otherwise show placeholder text."""
    if not os.path.exists(path):
        add_textbox(slide, left, top, width, Cm(1.5),
                    f"[图: {os.path.basename(path)}]", font_size=Pt(10), color=MID_GRAY)
        return None
    if height:
        pic = slide.shapes.add_picture(path, Emu(left), Emu(top), Emu(width), Emu(height))
    else:
        pic = slide.shapes.add_picture(path, Emu(left), Emu(top), width=Emu(width))
    if caption:
        add_textbox(slide, left, top + (height or Cm(7)) + Cm(0.1), width, Cm(0.5),
                    caption, font_size=Pt(9), color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    return pic

def add_round_box(slide, left, top, width, height, text, title=None, fill_color=None):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(0.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.3)
    tf.margin_right = Cm(0.3)
    tf.margin_top = Cm(0.2)

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        set_font(r, size=Pt(14), bold=True, color=CQU_BLUE)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = text
        set_font(r2, size=Pt(11), bold=False, color=DARK_GRAY)
    else:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        set_font(r, size=Pt(12), bold=False, color=DARK_GRAY)
    return shape

def add_metric_card(slide, left, top, width, height, label, value, color=CQU_BLUE):
    """Add a metric display card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)

    # Value
    add_textbox(slide, left + Cm(0.2), top + Cm(0.15), width - Cm(0.4), Cm(1.2),
                str(value), font_size=Pt(28), bold=True, color=color,
                alignment=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, left + Cm(0.2), top + height - Cm(1.0), width - Cm(0.4), Cm(0.8),
                label, font_size=Pt(10), bold=False, color=MID_GRAY,
                alignment=PP_ALIGN.CENTER)

def add_table_custom(slide, left, top, col_widths, headers, rows,
                     header_bg=TITLE_BG, header_fg=WHITE, font_size=Pt(10)):
    """Add a formatted table with header styling."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          Emu(left), Emu(top),
                                          Emu(total_w), Emu(Cm(0.7) * n_rows))
    table = table_shape.table

    for j, w in enumerate(col_widths):
        table.columns[j].width = Emu(w)

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        set_font(r, size=font_size, bold=True, color=header_fg)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        bg = WHITE if i % 2 == 0 else RGBColor(0xF5, 0xF7, 0xFA)
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val)
            set_font(r, size=font_size, bold=False, color=DARK_GRAY)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape

def add_bottom_bar(slide):
    """Add a thin bottom accent bar (CQU-style)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(SH - Cm(0.2)),
        SW, Emu(Cm(0.2))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CQU_BLUE
    shape.line.fill.background()

def add_page_number(slide, num, total):
    """Add page number at bottom right."""
    add_textbox(slide, Cm(22.5), Cm(17.5), Cm(3), Cm(0.5),
                f"{num} / {total}", font_size=Pt(8), color=MID_GRAY,
                alignment=PP_ALIGN.RIGHT)

# ======================== Main Build ========================

def build_ppt():
    global SW, SH
    import shutil, io, zipfile, os as _os
    from lxml import etree

    print("Creating clean presentation from template...")

    # Strategy: Copy the template and strip all slides at the ZIP level,
    # keeping only the slide masters, layouts, and theme.
    OUTPUT_TMP = os.path.join(os.path.dirname(OUTPUT), "_tmp_stripped.pptx")

    NSMAP = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'ct': 'http://schemas.openxmlformats.org/package/2006/content-types',
    }

    with zipfile.ZipFile(TEMPLATE, 'r') as zin:
        with zipfile.ZipFile(OUTPUT_TMP, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                # Skip all slide files (slide1.xml, slide2.xml, ...)
                if 'slides/slide' in item.filename:
                    continue
                # Skip slide relationships
                if 'slides/_rels/slide' in item.filename:
                    continue
                # Skip slide media references
                if 'slideMasters/_rels/slideMaster' in item.filename and 'slide' in item.filename:
                    pass  # keep master rels

                data = zin.read(item.filename)

                # For presentation.xml, strip all slide references from sldIdLst
                if item.filename == 'ppt/presentation.xml':
                    root = etree.fromstring(data)
                    sldIdLst = root.find('{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
                    if sldIdLst is not None:
                        for child in list(sldIdLst):
                            sldIdLst.remove(child)
                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

                # For [Content_Types].xml, remove slide content types
                if item.filename == '[Content_Types].xml':
                    root = etree.fromstring(data)
                    ns_ct = 'http://schemas.openxmlformats.org/package/2006/content-types'
                    for ov in root.findall(f'{{{ns_ct}}}Override'):
                        part = ov.get('PartName')
                        if part and '/slides/slide' in part:
                            root.remove(ov)
                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

                # For ppt/_rels/presentation.xml.rels, remove slide relationships
                if item.filename == 'ppt/_rels/presentation.xml.rels':
                    root = etree.fromstring(data)
                    for rel_elem in list(root):
                        target = rel_elem.get('Target', '')
                        if 'slides/slide' in target:
                            root.remove(rel_elem)
                    data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

                zout.writestr(item, data)

    # Now open the stripped presentation
    prs = Presentation(OUTPUT_TMP)
    SW = prs.slide_width
    SH = prs.slide_height
    print(f"Stripped presentation ready. Slide count: {len(prs.slides)}")

    TOTAL_SLIDES = 20

    # ====== Colors and constants ======
    CONTENT_LEFT = Cm(1.2)
    CONTENT_TOP = Cm(2.3)
    CONTENT_W = Cm(23.0)
    CONTENT_H = Cm(14.5)

    # ====== Slide 1: COVER ======
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片
    # Add title text
    add_textbox(slide, Cm(1.5), Cm(6.5), Cm(22), Cm(2.5),
                "空调用电负荷预测", font_size=Pt(42), bold=True, color=CQU_BLUE,
                alignment=PP_ALIGN.CENTER, font_name="微软雅黑")
    add_textbox(slide, Cm(1.5), Cm(9.0), Cm(22), Cm(1.5),
                "基于多算法对比的时间序列预测实验", font_size=Pt(22), bold=False,
                color=MID_GRAY, alignment=PP_ALIGN.CENTER, font_name="微软雅黑")
    add_textbox(slide, Cm(1.5), Cm(11.0), Cm(22), Cm(2.5),
                "重庆大学 大数据与软件学院 · 数据挖掘实验", font_size=Pt(16),
                bold=False, color=MID_GRAY, alignment=PP_ALIGN.CENTER, font_name="微软雅黑")
    add_textbox(slide, Cm(1.5), Cm(13.5), Cm(22), Cm(1.0),
                "汇报人：青课    日期：2026年6月", font_size=Pt(14),
                bold=False, color=MID_GRAY, alignment=PP_ALIGN.CENTER, font_name="微软雅黑")
    add_page_number(slide, 1, TOTAL_SLIDES)

    # ====== Slide 2: 目录 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # 自定义版式
    add_slide_title(slide, "目  录", "CONTENTS")
    add_bottom_bar(slide)

    toc_items = [
        ("01", "实验背景与目标", "基于UCI家庭用电数据集的多算法对比实验"),
        ("02", "数据预处理与特征工程", "缺失值处理、归一化、时间特征、滞后与滚动统计"),
        ("03", "模型构建与训练", "RF / XGBoost / LSTM / MLP-RF 四类算法实现"),
        ("04", "关键发现：数据泄漏", "从R²=0.999到发现物理公式泄漏的排查过程"),
        ("05", "实验结果与对比分析", "五模型性能指标、可视化诊断、特征重要性"),
        ("06", "实验总结与展望", "结论、启示、改进方向"),
    ]

    y_start = Cm(2.8)
    for i, (num, title, desc) in enumerate(toc_items):
        y = y_start + Cm(2.2) * i
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(Cm(2.0)), Emu(y), Emu(Cm(1.2)), Emu(Cm(1.2))
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = CQU_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        set_font(r, size=Pt(18), bold=True, color=WHITE)

        # Title
        add_textbox(slide, Cm(3.8), y + Cm(0.05), Cm(18), Cm(0.8),
                    title, font_size=Pt(22), bold=True, color=CQU_BLUE)
        # Description
        add_textbox(slide, Cm(3.8), y + Cm(0.75), Cm(18), Cm(0.5),
                    desc, font_size=Pt(12), bold=False, color=MID_GRAY)
    add_page_number(slide, 2, TOTAL_SLIDES)

    # ====== Slide 3: 实验背景与目标 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "01  实验背景与目标", "Background & Objectives")
    add_bottom_bar(slide)
    add_page_number(slide, 3, TOTAL_SLIDES)

    # Left column - 背景
    add_textbox(slide, Cm(1.5), Cm(2.5), Cm(10), Cm(0.8),
                "▎实验背景", font_size=Pt(20), bold=True, color=CQU_BLUE)
    bg_lines = [
        "• 电力负荷预测是智能电网和能源管理中的核心问题",
        "• 准确的短期负荷预测可优化发电调度、降低运营成本",
        "• 空调负荷在家庭用电中占比最高，季节性波动显著",
        "• 传统单一模型难以捕捉负荷的多尺度时序特征",
        "• 本实验使用UCI家庭用电数据集（47个月，200万+条记录）",
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(3.5), Cm(10.5), Cm(7),
                          bg_lines, font_size=Pt(13), color=DARK_GRAY)

    # Right column - 目标
    add_textbox(slide, Cm(13.0), Cm(2.5), Cm(10), Cm(0.8),
                "▎实验目标", font_size=Pt(20), bold=True, color=CQU_BLUE)
    obj_lines = [
        "① 实现完整的数据预处理与特征工程流水线",
        "② 构建并调优 4种不同范式的预测模型：",
        "    · 随机森林（Bagging集成）",
        "    · XGBoost（Boosting集成）",
        "    · LSTM（深度学习/循环神经网络）",
        "    · MLP-RF（神经网络+树模型融合）",
        "③ 使用5项指标（MSE/MAE/RMSE/MAPE/R²）全面评估",
        "④ 通过可视化和模型诊断进行深度对比分析",
        "⑤ 探究数据泄漏、特征重要性、融合策略等关键问题",
    ]
    add_multiline_textbox(slide, Cm(13.0), Cm(3.5), Cm(11.0), Cm(10),
                          obj_lines, font_size=Pt(12), color=DARK_GRAY)

    # Dataset info box at bottom
    add_round_box(slide, Cm(1.5), Cm(12.0), Cm(22.5), Cm(3.5),
                  "UCI Individual Household Electric Power Consumption：法国巴黎某家庭2006.12-2010.11的分钟级能耗数据，"
                  "含Global_active_power（目标）、Voltage、Sub_metering_1/2/3等7维变量。"
                  "逐小时重采样后34,168条记录，按70%/15%/15%时序分割为训练/验证/测试集。",
                  title="📊 数据集概况")

    # ====== Slide 4: 数据预处理全流程 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "02  数据预处理与特征工程", "Data Preprocessing & Feature Engineering")
    add_bottom_bar(slide)
    add_page_number(slide, 4, TOTAL_SLIDES)

    # Pipeline steps as horizontal flow
    steps = [
        ("原始数据加载", "127MB .txt →\npandas DataFrame\nfloat32内存优化"),
        ("缺失值处理", "删除1.25%\n含NaN行\n(~25K/2M条)"),
        ("逐小时重采样", "2,075,259行 →\n34,168行\n取每小时均值"),
        ("归一化", "Min-Max Scaling\n所有特征→[0,1]\n消除量纲差异"),
        ("特征构建", "时间循环编码 +\n滞后特征 +\n滚动窗口统计"),
        ("数据分割", "时序分割70/15/15\n防止未来信息\n泄漏到训练集"),
    ]

    x_pos = Cm(0.8)
    for i, (title, desc) in enumerate(steps):
        # Step box
        add_round_box(slide, x_pos, Cm(2.8), Cm(3.8), Cm(4.5), desc, title=title,
                      fill_color=RGBColor(0xEE, 0xF2, 0xF9))
        # Arrow between steps
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Emu(x_pos + Cm(4.0)), Emu(Cm(4.5)),
                Emu(Cm(0.8)), Emu(Cm(0.6))
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = CQU_RED
            arrow.line.fill.background()
        x_pos += Cm(4.0)

    # Key feature details below
    left_col_lines = [
        ("【时间循环编码】", True, CQU_BLUE, Pt(13)),
        "• hour_sin/cos：sin(2π·hour/24), cos(2π·hour/24)",
        "  → 保留24小时周期的循环性质（23时与0时相邻）",
        "• month_sin/cos：sin(2π·month/12), cos(2π·month/12)",
        "  → 保留12个月的年周期循环",
        "• dayofweek, is_weekend：工作日/周末二元标识",
        ("【滞后特征 (6个)】", True, CQU_BLUE, Pt(13)),
        "• lag1h, lag2h, lag3h, lag6h, lag12h, lag24h",
        "• 编码时序自相关：前一时刻负荷是最直接的预测信号",
    ]
    add_multiline_textbox(slide, Cm(1.2), Cm(8.2), Cm(11.0), Cm(8.0),
                          left_col_lines, font_size=Pt(11), color=DARK_GRAY,
                          line_spacing=1.3)

    right_col_lines = [
        ("【滚动窗口统计 (12个)】", True, CQU_BLUE, Pt(13)),
        "• 窗口大小：6小时 / 12小时 / 24小时",
        "• 每个窗口计算：mean, std, min, max",
        "• 例：roll_mean_24h = 过去24小时负荷均值",
        "• 捕捉局部趋势、波动性和极值特征",
        ("【特征汇总】", True, CQU_BLUE, Pt(13)),
        "• 原始特征：5维（排除泄漏特征后）",
        "• 时间特征：8维（hour/dayofweek/month/is_weekend + 4循环编码）",
        "• 滞后特征：6维",
        "• 滚动统计：12维",
        "• 总计：30维特征向量（去除NaN后34,144行）",
    ]
    add_multiline_textbox(slide, Cm(13.0), Cm(8.2), Cm(11.5), Cm(8.0),
                          right_col_lines, font_size=Pt(11), color=DARK_GRAY,
                          line_spacing=1.3)

    # ====== Slide 5: EDA 探索性数据分析 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "02  探索性数据分析（EDA）", "Exploratory Data Analysis")
    add_bottom_bar(slide)
    add_page_number(slide, 5, TOTAL_SLIDES)

    add_image_safe(slide, os.path.join(FIGURES, "eda_correlation_heatmap.png"),
                   Cm(0.8), Cm(2.2), Cm(10.5), Cm(8.5),
                   caption="图1: 变量相关性热力图")
    add_image_safe(slide, os.path.join(FIGURES, "eda_seasonal_decompose.png"),
                   Cm(11.8), Cm(2.2), Cm(13.0), Cm(8.5),
                   caption="图2: 季节性分解（Trend + Seasonal + Residual）")

    # Insights
    insights = [
        ("关键发现：", True, CQU_RED, Pt(12)),
        "• Global_active_power 与 Global_intensity 相关系数≈1.0",
        "   → 物理公式 P=V×I 决定了此关系 → 必须排除（数据泄漏）",
        "• Sub_metering_3（空调/热水器）与目标相关性最高（~0.7）",
        "   → 空调是家庭负荷的主要驱动力",
        "• 季节性分解显示：周末用电高于工作日；存在年周期趋势",
        "• 功率热力图显示：白天10:00-18:00用电显著高于夜间",
    ]
    add_multiline_textbox(slide, Cm(1.2), Cm(11.5), Cm(23.0), Cm(5.0),
                          insights, font_size=Pt(10), color=DARK_GRAY,
                          line_spacing=1.25)

    # ====== Slide 6: 模型架构总览 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "03  模型构建总览", "Model Architecture Overview")
    add_bottom_bar(slide)
    add_page_number(slide, 6, TOTAL_SLIDES)

    models = [
        ("随机森林\nRandom Forest", "Bagging集成\n200棵树, max_depth=25\nmin_samples_split=5\n抗噪、不易过拟合",
         CQU_BLUE),
        ("XGBoost", "Boosting集成\n400棵树, max_depth=6\nlr=0.05, 正则化\n高效、自动处理缺失值",
         ACCENT_BLUE),
        ("LSTM", "3层循环神经网络\nhidden_dim=128\ndropout=0.3\n序列建模、记忆效应",
         ACCENT_GREEN),
        ("MLP-RF融合", "MLP(3层256→128→64)\n+ RF简单平均融合\nBatchNorm+Dropout\n错误模式互补",
         ACCENT_RED),
    ]

    for i, (name, desc, color) in enumerate(models):
        x = Cm(0.8) + Cm(6.2) * i
        # Model card
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(x), Emu(Cm(2.8)), Emu(Cm(5.8)), Emu(Cm(10.0))
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        # Header with color
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(x), Emu(Cm(2.8)), Emu(Cm(5.8)), Emu(Cm(2.2))
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = color
        header_bar.line.fill.background()
        add_textbox(slide, x + Cm(0.2), Cm(3.0), Cm(5.4), Cm(1.8),
                    name, font_size=Pt(18), bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_multiline_textbox(slide, x + Cm(0.3), Cm(5.5), Cm(5.2), Cm(7.0),
                              desc.split("\n"), font_size=Pt(12), color=DARK_GRAY,
                              line_spacing=1.5)

    # Training setup at bottom
    add_textbox(slide, Cm(1.5), Cm(13.5), Cm(22), Cm(3.0),
                "训练环境：PyTorch 2.12 + CUDA 13.0 + NVIDIA RTX 4060 Laptop (8GB)  |  "
                "超参搜索：RandomizedSearchCV / 3折交叉验证  |  "
                "早停策略：patience=15  |  "
                "评估指标：MSE / MAE / RMSE / MAPE / R²",
                font_size=Pt(11), bold=False, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ====== Slide 7: 随机森林 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "03  模型详解：随机森林  Random Forest", "Model Details: RF")
    add_bottom_bar(slide)
    add_page_number(slide, 7, TOTAL_SLIDES)

    rf_left = [
        ("▎算法原理", True, CQU_BLUE, Pt(18)),
        "",
        "• 基于Bagging（Bootstrap Aggregating）的集成方法",
        "• 并行构建多个决策树，最终取所有树的预测均值",
        "• 两个随机性来源：",
        "  ① 样本随机：每棵树从训练集中有放回抽样",
        "  ② 特征随机：每次分裂只考虑部分特征子集",
        "• 通过增加树的多样性来降低方差，天然抗过拟合",
        "",
        ("▎本次实验做了什么", True, CQU_BLUE, Pt(18)),
        "",
        "• 使用 scikit-learn RandomForestRegressor",
        "• RandomizedSearchCV 超参搜索：15轮 × 3折CV",
        "• 搜索空间：n_estimators [100,200,300],",
        "  max_depth [10,15,20,25,None],",
        "  min_samples_split [2,5,10],",
        "  min_samples_leaf [1,2,4],",
        "  max_features [0.5, 0.7, 1.0]",
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(2.5), Cm(12.5), Cm(14.0),
                          rf_left, font_size=Pt(12), color=DARK_GRAY)

    rf_right = [
        ("▎最优参数", True, CQU_BLUE, Pt(18)),
        "",
        "n_estimators = 200",
        "max_depth = 25",
        "min_samples_split = 5",
        "min_samples_leaf = 4",
        "max_features = 0.5",
        "",
        ("▎核心优势", True, CQU_BLUE, Pt(18)),
        "",
        "✓ 无需特征归一化（树模型天然不受量纲影响）",
        "✓ 可输出特征重要性，帮助理解数据",
        "✓ 对异常值不敏感",
        "✓ 训练可高度并行化",
        "",
        ("训练用时：1201s（单线程）", False, MID_GRAY, Pt(12)),
        ("数据量：29,022条 × 30维特征", False, MID_GRAY, Pt(12)),
    ]
    add_multiline_textbox(slide, Cm(15.0), Cm(2.5), Cm(9.5), Cm(14.0),
                          rf_right, font_size=Pt(13), color=DARK_GRAY)

    # ====== Slide 8: XGBoost ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "03  模型详解：XGBoost", "Model Details: XGBoost")
    add_bottom_bar(slide)
    add_page_number(slide, 8, TOTAL_SLIDES)

    xgb_left = [
        ("▎算法原理", True, CQU_BLUE, Pt(18)),
        "",
        "• 基于Gradient Boosting框架，树串行构建",
        "• 每棵新树拟合前一棵树的残差（梯度方向）",
        "• 目标函数 = 损失函数 + 正则化项（Ω）",
        "• 正则化：控制叶子节点数量与权重，防止过拟合",
        "• tree_method='hist'：基于直方图的近似分裂算法",
        "  → 大幅加速训练，适合大规模数据",
        "• 内置缺失值自动学习：无需手动填充",
        "• 反向剪枝（后剪枝）：先建树再剪枝，避免局部最优",
        "",
        ("▎本次实验做了什么", True, CQU_BLUE, Pt(18)),
        "",
        "• 使用 xgboost XGBRegressor (tree_method='hist')",
        "• RandomizedSearchCV 超参搜索：20轮 × 3折CV",
        "• 搜索空间：n_estimators [100-600],",
        "  max_depth [4-12], lr [0.01-0.1],",
        "  subsample/colsample [0.7-1.0],",
        "  正则化 reg_alpha/reg_lambda",
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(2.5), Cm(12.5), Cm(14.0),
                          xgb_left, font_size=Pt(12), color=DARK_GRAY)

    xgb_right = [
        ("▎最优参数", True, CQU_BLUE, Pt(18)),
        "",
        "n_estimators = 400",
        "max_depth = 6",
        "learning_rate = 0.05",
        "subsample = 0.8",
        "colsample_bytree = 0.7",
        "reg_alpha = 0,  reg_lambda = 1.5",
        "",
        ("▎核心优势", True, CQU_BLUE, Pt(18)),
        "",
        "✓ 训练速度极快（106s vs RF的1201s）",
        "✓ L1/L2正则化天然防过拟合",
        "✓ 内置缺失值处理，工程实用性最强",
        "✓ 实际部署场景的最佳选择",
        "",
        ("训练用时：106s", False, MID_GRAY, Pt(12)),
        ("工业界时序预测的SOTA基线", False, MID_GRAY, Pt(12)),
    ]
    add_multiline_textbox(slide, Cm(15.0), Cm(2.5), Cm(9.5), Cm(14.0),
                          xgb_right, font_size=Pt(13), color=DARK_GRAY)

    # ====== Slide 9: LSTM ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "03  模型详解：LSTM 长短期记忆网络", "Model Details: LSTM")
    add_bottom_bar(slide)
    add_page_number(slide, 9, TOTAL_SLIDES)

    lstm_left = [
        ("▎算法原理", True, CQU_BLUE, Pt(18)),
        "",
        "• RNN的变体，解决长序列梯度消失/爆炸问题",
        "• 三个门控机制：",
        "  · 遗忘门 f_t = σ(W_f·[h_{t-1}, x_t] + b_f)",
        "    决定丢弃哪些旧信息",
        "  · 输入门 i_t = σ(W_i·[h_{t-1}, x_t] + b_i)",
        "    决定将哪些新信息存入细胞状态",
        "  · 输出门 o_t = σ(W_o·[h_{t-1}, x_t] + b_o)",
        "    决定输出哪些信息",
        "• 细胞状态 C_t 作为长期记忆通路",
        "",
        ("▎本次实验做了什么", True, CQU_BLUE, Pt(18)),
        "",
        "• PyTorch构建3层LSTM (hidden_dim=128, dropout=0.3)",
        "• 输入：过去24小时 × 10维多变量序列",
        "• 输出：下一小时负荷预测值",
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(2.5), Cm(12.5), Cm(14.0),
                          lstm_left, font_size=Pt(12), color=DARK_GRAY)

    lstm_right = [
        ("▎训练配置", True, CQU_BLUE, Pt(18)),
        "",
        "优化器：Adam (lr=1e-3)",
        "学习率衰减：ReduceLROnPlateau",
        "梯度裁剪：max_norm=1.0",
        "早停：patience=15轮",
        "Batch Size：64",
        "参数量：~200K",
        "",
        ("▎训练曲线", True, CQU_BLUE, Pt(18)),
        "",
        "• 约第48轮达到最优验证loss",
        "• 第63轮触发早停",
        "• 训练loss持续下降但验证loss在~0.007处停滞",
        "• 模型已达到容量上限",
    ]
    add_multiline_textbox(slide, Cm(15.0), Cm(2.5), Cm(9.5), Cm(8.0),
                          lstm_right, font_size=Pt(13), color=DARK_GRAY)

    add_image_safe(slide, os.path.join(FIGURES, "lstm_training_curves.png"),
                   Cm(15.0), Cm(10.5), Cm(9.5), Cm(6.0),
                   caption="图: LSTM训练与验证损失曲线")

    # ====== Slide 10: MLP-RF融合 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "03  模型详解：MLP-RF 融合模型", "Model Details: MLP-RF Ensemble")
    add_bottom_bar(slide)
    add_page_number(slide, 10, TOTAL_SLIDES)

    mlp_left = [
        ("▎融合策略对比", True, CQU_BLUE, Pt(18)),
        "",
        ("策略1：Ridge Stacking", True, CQU_RED, Pt(14)),
        "• 以RF和MLP预测值为特征，Ridge回归作元学习器",
        "• 结果：Ridge给MLP的权重≈-0.07（几乎为零）",
        "• 融合退化 ≈ 单独使用RF（R²=0.8905）",
        "• 失败原因：两个模型预测高度相关(ρ>0.95)",
        "  线性元学习器无法从中提取额外信息",
        "",
        ("策略2：简单算术平均 ✓（最终采用）", True, ACCENT_GREEN, Pt(14)),
        "• ensemble_pred = (RF_pred + MLP_pred) / 2",
        "• 强制各50%权重，偏差模式部分抵消",
        "• R² = 0.9007（优于所有单一模型）",
        "",
        "启示：基模型预测高度相关时，简单平均",
        "      可能比复杂元学习更有效。",
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(2.5), Cm(12.5), Cm(14.0),
                          mlp_left, font_size=Pt(12), color=DARK_GRAY)

    mlp_right = [
        ("▎MLP网络结构", True, CQU_BLUE, Pt(18)),
        "",
        "输入层: 30维特征",
        "  ↓",
        "Linear(30→256) + BatchNorm + ReLU + Dropout(0.3)",
        "  ↓",
        "Linear(256→128) + BatchNorm + ReLU + Dropout(0.3)",
        "  ↓",
        "Linear(128→64) + BatchNorm + ReLU + Dropout(0.3)",
        "  ↓",
        "Linear(64→1)  → 输出预测值",
        "",
        ("▎为什么MLP能接近树模型？", True, CQU_BLUE, Pt(14)),
        "",
        "• BatchNorm：加速收敛、稳定训练",
        "• Dropout(0.3)：防止过拟合",
        "• 30维显式特征已编码时序依赖",
        "  → 降低了对模型隐式学习能力的要求",
        "• 训练仅23s，与树模型性能接近",
    ]
    add_multiline_textbox(slide, Cm(15.0), Cm(2.5), Cm(9.5), Cm(14.0),
                          mlp_right, font_size=Pt(12), color=DARK_GRAY)

    # ====== Slide 11: 数据泄漏 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "04  关键发现：数据泄漏排查", "Critical Discovery: Data Leakage Detection")
    add_bottom_bar(slide)
    add_page_number(slide, 11, TOTAL_SLIDES)

    # Big red alert box
    alert_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(Cm(1.5)), Emu(Cm(2.5)), Emu(Cm(22.5)), Emu(Cm(4.0))
    )
    alert_shape.fill.solid()
    alert_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
    alert_shape.line.color.rgb = CQU_RED
    alert_shape.line.width = Pt(2)

    alert_lines = [
        ("⚠️  实验初期，模型R²一度达到 0.999", True, CQU_RED, Pt(22)),
        ("表面近乎完美 —— 但这不是好消息，而是警报。", False, DARK_GRAY, Pt(14)),
        ("特征重要性分析揭示：Global_intensity（电流）占据 99.9% 的重要性。", False, DARK_GRAY, Pt(14)),
    ]
    add_multiline_textbox(slide, Cm(2.0), Cm(2.8), Cm(21.5), Cm(3.5),
                          alert_lines, font_size=Pt(14), line_spacing=1.5)

    # Explanation
    explain_lines = [
        ("▎根因分析", True, CQU_BLUE, Pt(18)),
        "",
        "物理定律：P = V × I （功率 = 电压 × 电流）",
        "在法国电网中，电压 V≈240V（几乎恒定），因此：",
        "",
        "    Global_active_power ≈ 240 × Global_intensity",
        "",
        "将Global_intensity作为特征 → 模型只需学会乘以240即得到答案",
        "这不是「预测」，这是「读取答案」——典型的数据泄漏（Data Leakage）。",
        "",
        ("▎修复措施", True, CQU_BLUE, Pt(18)),
        "• 排除特征：Global_intensity, Global_reactive_power",
        "• 仅保留：Voltage, Sub_metering_1, Sub_metering_2, Sub_metering_3",
        "• 修复后，R²从0.999降至0.89~0.90 → 这才是真实的预测性能",
        "",
        ("💡 教训：如果某个特征Importance > 90%，立即审视其与目标变量的物理/逻辑关系。", True, CQU_RED, Pt(14)),
        ("     R²=0.999不是庆祝的理由，而是检查数据泄漏的警报。", True, CQU_RED, Pt(14)),
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(7.0), Cm(22.5), Cm(10.0),
                          explain_lines, font_size=Pt(12), color=DARK_GRAY,
                          line_spacing=1.2)

    # ====== Slide 12: 实验结果 —— 性能总表 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "05  实验结果：五模型性能对比", "Results: Model Performance Comparison")
    add_bottom_bar(slide)
    add_page_number(slide, 12, TOTAL_SLIDES)

    # Metrics table
    col_w = [Cm(4.0), Cm(3.5), Cm(3.5), Cm(3.5), Cm(3.5), Cm(4.0)]
    headers = ["模型 Model", "MSE ↓", "MAE ↓", "RMSE ↓", "MAPE(%) ↓", "R² ↑"]
    rows = [
        ["MLP-RF Avg", "0.0012", "0.0228", "0.0347", "24.60", "0.9007"],
        ["XGBoost",    "0.0013", "0.0232", "0.0358", "23.79", "0.8947"],
        ["Random Forest","0.0013","0.0238", "0.0359", "25.09", "0.8937"],
        ["MLP",         "0.0013", "0.0240", "0.0365", "26.15", "0.8904"],
        ["LSTM",        "0.0051", "0.0489", "0.0715", "53.92", "0.5793"],
    ]
    add_table_custom(slide, Cm(1.5), Cm(2.5), col_w, headers, rows)

    # Key findings
    findings = [
        ("🏆 最佳R²：MLP-RF Avg（0.9007）— 简单平均融合超越所有单一模型", False, ACCENT_BLUE, Pt(14)),
        ("⚡ 最佳MAPE：XGBoost（23.79%）— 实际部署场景的最优选择", False, ACCENT_BLUE, Pt(14)),
        ("📉 LSTM表现最差（MAPE=53.9%, R²=0.579）— 特征维度与数据量不足", False, ACCENT_RED, Pt(14)),
        ("📊 树模型（RF/XGBoost）表现稳健 — 丰富显式特征工程是关键", False, DARK_GRAY, Pt(14)),
        ("🧠 MLP（3层）训练仅23s，性能接近树模型 — NN在表格化时序上同样有竞争力", False, DARK_GRAY, Pt(14)),
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(7.8), Cm(22.5), Cm(4.0),
                          findings, font_size=Pt(11), color=DARK_GRAY, line_spacing=1.5)

    # Comparison chart
    add_image_safe(slide, os.path.join(FIGURES, "final_summary.png"),
                   Cm(1.5), Cm(11.0), Cm(22.5), Cm(6.0),
                   caption="图: 五模型综合对比（MAPE/R²/RMSE柱状图 + 预测曲线 + 关键发现）")

    # ====== Slide 13: 预测曲线 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "05  预测结果可视化", "Prediction Visualization")
    add_bottom_bar(slide)
    add_page_number(slide, 13, TOTAL_SLIDES)

    add_image_safe(slide, os.path.join(FIGURES, "all_models_7days.png"),
                   Cm(0.5), Cm(2.2), Cm(24.5), Cm(7.0),
                   caption="图: 五模型预测曲线对比（测试集前7天 / 168小时）")
    add_image_safe(slide, os.path.join(FIGURES, "all_models_scatter.png"),
                   Cm(0.5), Cm(9.5), Cm(24.5), Cm(7.0),
                   caption="图: 五模型预测值 vs. 真实值散点图（越接近对角线 y=x 表示预测越准确）")

    # ====== Slide 14: 特征重要性 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "05  特征重要性分析", "Feature Importance Analysis")
    add_bottom_bar(slide)
    add_page_number(slide, 14, TOTAL_SLIDES)

    fi_insights = [
        ("▎RF与XGBoost特征重要性排序高度一致：", True, CQU_BLUE, Pt(16)),
        "",
        "1. Sub_metering_3（空调/热水器） — 重要性 ~38%",
        "   → 空调是家庭电力负荷的最强驱动因素",
        "",
        "2. lag1h（前1小时负荷） — 重要性 ~12-21%",
        "   → 电力负荷具有强时序自相关性",
        "   → 最近时刻的负荷是预测下一时刻最直接的信息",
        "",
        "3. Sub_metering_1 + Sub_metering_2（厨房+洗衣间） — 合计 ~20%",
        "   → 特定电器的使用模式提供了额外的预测信号",
        "",
        "4. 时间循环特征（hour_sin/cos, month_cos） — 合计 ~5%",
        "   → 日周期和年周期特征为模型提供了时间上下文",
        "",
        "5. Voltage（电压） — ~2-3%",
        "   → 电网电压的微小波动也包含一定预测信息",
        "",
        ("💡 结论：显式特征工程（滞后+滚动统计+循环编码）直接编码了时序依赖，", False, DARK_GRAY, Pt(13)),
        ("     大幅降低了模型的学习难度。这是树模型在本任务中表现出色的关键原因。", False, DARK_GRAY, Pt(13)),
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(2.5), Cm(22.5), Cm(14.0),
                          fi_insights, font_size=Pt(12), color=DARK_GRAY, line_spacing=1.15)

    # ====== Slide 15: LSTM 为什么差 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "05  深度分析：LSTM为何表现不佳？", "Analysis: Why LSTM Underperformed")
    add_bottom_bar(slide)
    add_page_number(slide, 15, TOTAL_SLIDES)

    lstm_analysis = [
        ("LSTM 在本次实验中表现最差（MAPE=53.9%, R²=0.579），三个主要原因：", False, CQU_RED, Pt(16)),
        "",
        ("❶ 特征不对称（最核心原因）", True, CQU_BLUE, Pt(16)),
        "• LSTM：仅使用10维原始特征序列（含自回归项）",
        "• RF/XGBoost：获得30维显式工程特征（滞后+滚动统计+循环编码）",
        "• LSTM需要从序列中「隐式」学习这些模式",
        "  → 对模型容量和数据量的要求远超树模型",
        "• 这是实验设计上的不对称：给树模型做了充分特征工程，",
        "  但没有给LSTM平等的特征支持",
        "",
        ("❷ 数据量不足", True, CQU_BLUE, Pt(16)),
        "• ~29K训练样本 对 3层LSTM（~200K参数）偏少",
        "• 树模型在小样本场景下天然具有先验优势",
        "• 深度学习通常需要更大规模数据才能发挥优势",
        "",
        ("❸ 任务特性", True, CQU_BLUE, Pt(16)),
        "• 对于已做好特征工程的表格化时序预测",
        "  梯度提升树（XGBoost）通常是更优选择",
        "• LSTM的优势场景：原始长序列、强非线性多变量交互、",
        "  大数据量下端到端学习",
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(2.3), Cm(22.5), Cm(15.0),
                          lstm_analysis, font_size=Pt(11), color=DARK_GRAY,
                          line_spacing=1.15)

    # ====== Slide 16: EDA更多图表 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "05  更多可视化分析", "Additional Visualizations")
    add_bottom_bar(slide)
    add_page_number(slide, 16, TOTAL_SLIDES)

    add_image_safe(slide, os.path.join(FIGURES, "eda_power_heatmap.png"),
                   Cm(0.5), Cm(2.2), Cm(12.0), Cm(6.5),
                   caption="图: 功率变化热力图（日期×小时，最后60天）")
    add_image_safe(slide, os.path.join(FIGURES, "rf_xgb_residuals.png"),
                   Cm(13.0), Cm(2.2), Cm(12.0), Cm(6.5),
                   caption="图: RF与XGBoost残差分布（近似正态，无系统性偏差）")
    add_image_safe(slide, os.path.join(FIGURES, "eda_acf_pacf.png"),
                   Cm(0.5), Cm(9.2), Cm(12.0), Cm(6.5),
                   caption="图: ACF/PACF自相关/偏自相关函数图")
    add_image_safe(slide, os.path.join(FIGURES, "eda_distributions.png"),
                   Cm(13.0), Cm(9.2), Cm(12.0), Cm(6.5),
                   caption="图: 各变量分布直方图")

    # ====== Slide 17: Stacking vs Average ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "05  融合策略对比：Stacking vs. 简单平均", "Ensemble Strategy Comparison")
    add_bottom_bar(slide)
    add_page_number(slide, 17, TOTAL_SLIDES)

    compare_lines = [
        ("实验对比了两种MLP-RF融合策略：", False, DARK_GRAY, Pt(15)),
        "",
        ("▎Ridge Stacking → R² = 0.8905（不如单独RF的0.8937）", True, CQU_RED, Pt(16)),
        "• Ridge元学习器给MLP的权重≈ -0.07（几乎为零）",
        "• 融合退化为单独使用RF，完全失去了集成的意义",
        "• 失败原因：RF和MLP在测试集上的预测相关系数 ρ > 0.95",
        "  → 两个模型输出高度共线，线性元学习器无法从中提取新信息",
        "",
        ("▎简单算术平均 → R² = 0.9007 ✨（最优）", True, ACCENT_GREEN, Pt(16)),
        "• ensemble_pred = (RF_pred + MLP_pred) / 2",
        "• 强制给每个模型50%权重",
        "• 两个模型偏差模式不完全重合，50/50平均恰好使误差部分抵消",
        "• 实现了真正的「集成学习」——整体优于部分",
        "",
        ("💡 核心启示：", True, CQU_BLUE, Pt(16)),
        "当基学习器的预测高度相关时，复杂的元学习（Stacking）可能事与愿违。",
        "简单平均反而因强制均分权重而获得更好的泛化效果。",
        "在集成学习中，「简单」有时候确实比「复杂」更有效。",
    ]
    add_multiline_textbox(slide, Cm(1.5), Cm(2.5), Cm(22.5), Cm(14.5),
                          compare_lines, font_size=Pt(12), color=DARK_GRAY,
                          line_spacing=1.2)

    # ====== Slide 18: 实验总结 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "06  实验总结", "Summary & Conclusions")
    add_bottom_bar(slide)
    add_page_number(slide, 18, TOTAL_SLIDES)

    summary_items = [
        ("1", "XGBoost — 实际部署的最佳选择",
         "MAPE=23.8%, 训练106s, 精度与效率的最佳平衡点",
         ACCENT_BLUE),
        ("2", "MLP-RF融合 — 最高R² (0.9007)",
         "简单平均超越所有单一模型, 证明了模型集成的价值",
         ACCENT_GREEN),
        ("3", "数据泄漏 — 最有价值的教训",
         "R²=0.999不是庆祝而是警报, 物理理解比高指标更重要",
         CQU_RED),
        ("4", "特征工程 — 性能上限的决定因素",
         "Sub_metering_3(~38%) + lag1h(~21%) = 最强预测信号",
         CQU_BLUE),
        ("5", "LSTM — 深度学习并非万能",
         "受限于特征维度和数据量, 提醒我们选对工具比选「高级」工具更重要",
         MID_GRAY),
        ("6", "简单平均 > Stacking",
         "当基学习器预测高度相关时, 简单的集成策略可能更有效",
         MID_GRAY),
    ]

    for i, (num, title, desc, color) in enumerate(summary_items):
        y = Cm(2.5) + Cm(2.5) * i
        # Number
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(Cm(1.5)), Emu(y + Cm(0.2)),
            Emu(Cm(0.9)), Emu(Cm(0.9))
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        set_font(r, size=Pt(14), bold=True, color=WHITE)

        # Title
        add_textbox(slide, Cm(2.8), y, Cm(20), Cm(0.7),
                    title, font_size=Pt(18), bold=True, color=color)
        # Description
        add_textbox(slide, Cm(2.8), y + Cm(0.7), Cm(20), Cm(0.6),
                    desc, font_size=Pt(12), bold=False, color=MID_GRAY)

    # ====== Slide 19: 展望 ======
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_slide_title(slide, "06  不足与展望", "Limitations & Future Work")
    add_bottom_bar(slide)
    add_page_number(slide, 19, TOTAL_SLIDES)

    future_items = [
        ("当前不足", [
            "• LSTM特征不对称：未平等对待深度模型",
            "• 样本量有限：~29K样本对深度模型偏少",
            "• 超参搜索粗糙：仅15-20轮RandomizedSearch",
            "• 未引入外部特征：如气象数据（温度/湿度）",
            "• 模型单一：未尝试Transformer、Informer等新架构",
        ]),
        ("改进方向", [
            "• 公平实验：为所有模型提供相同的特征集",
            "• 引入气象数据：温度、湿度是空调负荷的核心驱动",
            "• 数据扩充：使用更长时间跨度的数据集",
            "• 尝试新架构：Transformer/Informer/Autoformer",
            "• 更系统的超参优化：Optuna/Bayesian Optimization",
            "• 多步预测：从1步→24步（真正的日负荷预测）",
            "• 不确定性量化：预测区间而不仅是点估计",
        ]),
    ]

    for col_idx, (title, items) in enumerate(future_items):
        x = Cm(1.5) + Cm(12.5) * col_idx
        add_textbox(slide, x, Cm(2.8), Cm(11), Cm(0.8),
                    f"▎{title}", font_size=Pt(20), bold=True, color=CQU_BLUE)
        add_multiline_textbox(slide, x, Cm(3.8), Cm(11), Cm(12.5),
                              items, font_size=Pt(13), color=DARK_GRAY,
                              line_spacing=1.5)

    # ====== Slide 20: 结束页 ======
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片 (same as cover)
    add_textbox(slide, Cm(1.5), Cm(5.0), Cm(22), Cm(2.5),
                "感谢聆听！", font_size=Pt(48), bold=True, color=CQU_BLUE,
                alignment=PP_ALIGN.CENTER, font_name="微软雅黑")
    add_textbox(slide, Cm(1.5), Cm(8.0), Cm(22), Cm(1.5),
                "Thank You", font_size=Pt(28), bold=False,
                color=MID_GRAY, alignment=PP_ALIGN.CENTER, font_name="微软雅黑")
    add_textbox(slide, Cm(1.5), Cm(10.5), Cm(22), Cm(2.0),
                "空调用电负荷预测 — 基于多算法对比的时间序列预测实验",
                font_size=Pt(16), bold=False, color=MID_GRAY,
                alignment=PP_ALIGN.CENTER, font_name="微软雅黑")
    add_textbox(slide, Cm(1.5), Cm(12.5), Cm(22), Cm(1.5),
                "汇报人：青课  |  重庆大学 大数据与软件学院  |  数据挖掘实验",
                font_size=Pt(14), bold=False, color=MID_GRAY,
                alignment=PP_ALIGN.CENTER, font_name="微软雅黑")
    add_textbox(slide, Cm(1.5), Cm(14.5), Cm(22), Cm(1.0),
                "欢迎提问与交流", font_size=Pt(18), bold=True, color=CQU_RED,
                alignment=PP_ALIGN.CENTER, font_name="微软雅黑")
    add_page_number(slide, 20, TOTAL_SLIDES)

    # ======================== Save ========================
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"\n✅ PPT saved to: {OUTPUT}")
    print(f"   Total slides: {TOTAL_SLIDES}")

if __name__ == "__main__":
    build_ppt()
