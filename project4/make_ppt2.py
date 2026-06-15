"""
Generate a brand-new defense PPT for Project 4: Air-Conditioning Power Load Forecasting
=========================================================================
No external template — completely self-contained design.
Font-scaled for presentation readability with image annotations.
"""
import os
from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(BASE, "output", "空调负荷预测_答辩汇报_新版.pptx")
FIGURES = os.path.join(BASE, "output")

# ── Colors ──────────────────────────────────────────
C_DARK  = RGBColor(0x1B, 0x2A, 0x4A)
C_BLUE  = RGBColor(0x2D, 0x5F, 0x8A)
C_ACC   = RGBColor(0xE8, 0x6A, 0x17)
C_RED   = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_LGRAY = RGBColor(0xAA, 0xAA, 0xAA)
C_BG    = RGBColor(0xF7, 0xF8, 0xFC)
C_CARD  = RGBColor(0xFF, 0xFF, 0xFF)
C_DIV   = RGBColor(0xE0, 0xE3, 0xEB)
C_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
C_SUB    = RGBColor(0xB0, 0xBD, 0xD3)

SW = Cm(33.867)
SH = Cm(19.05)

# ═══════════════════ Helpers ═════════════════════════
def _ea(rPr, fn):
    for c in list(rPr):
        if c.tag == qn('a:ea'): rPr.remove(c)
    ea = rPr.makeelement(qn('a:ea'), {}); ea.set('typeface', fn); rPr.append(ea)

def font(run, size=Pt(16), bold=False, color=C_DARK, name="微软雅黑"):
    run.font.size = size; run.font.name = name; run.font.bold = bold
    if color: run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr(); _ea(rPr, name)
    lat = rPr.makeelement(qn('a:latin'), {}); lat.set('typeface', name)
    for c in list(rPr):
        if c.tag == qn('a:latin'): rPr.remove(c)
    rPr.append(lat)

def txt(s, l, t, w, h, text, size=Pt(16), bold=False, color=C_DARK, align=PP_ALIGN.LEFT, name="微软雅黑", sp=1.2):
    tb = s.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = Pt(size.pt * sp)
    r = p.add_run(); r.text = text; font(r, size, bold, color, name)
    return tb

def mtxt(s, l, t, w, h, lines, size=Pt(15), color=C_DARK, align=PP_ALIGN.LEFT, name="微软雅黑", sp=1.25):
    tb = s.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, li in enumerate(lines):
        if isinstance(li, str): text, b, c, sz = li, False, color, size
        else:
            text = li[0]; b = li[1] if len(li)>1 else False
            c = li[2] if len(li)>2 else color; sz = li[3] if len(li)>3 else size
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = Pt(sz.pt * sp)
        r = p.add_run(); r.text = text; font(r, sz, b, c, name)
    return tb

def rr(s, l, t, w, h, fill=C_CARD, border=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border: shp.line.color.rgb = border; shp.line.width = Pt(1)
    else: shp.line.fill.background()
    return shp

def oval(s, l, t, w, h, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(l), Emu(t), Emu(w), Emu(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()
    return shp

def img(s, path, l, t, w, h):
    if not os.path.exists(path):
        txt(s, l, t, w, Cm(1), f"[缺图]", Pt(10), color=C_LGRAY); return None
    return s.shapes.add_picture(path, Emu(l), Emu(t), Emu(w), Emu(h))

def page_bar(s, num, title, subtitle=None):
    rr(s, Cm(0), Cm(0), SW, Cm(2.8), C_DARK)
    rr(s, Cm(0), Cm(2.8), SW, Cm(0.08), C_ACC)
    txt(s, Cm(2.0), Cm(0.3), Cm(28), Cm(1.5), title, Pt(32), True, C_WHITE)
    if subtitle:
        txt(s, Cm(2.0), Cm(1.7), Cm(28), Cm(0.6), subtitle, Pt(15), color=C_SUB)
    txt(s, Cm(30), Cm(17.8), Cm(3), Cm(0.5), str(num), Pt(14), color=C_LGRAY, align=PP_ALIGN.RIGHT)

def add_bottom(s):
    rr(s, Cm(0), SH - Cm(0.15), SW, Cm(0.15), C_DARK)

def card(s, l, t, w, h, title, body_lines, tc=C_BLUE, bs=Pt(14)):
    rr(s, l, t, w, h, C_CARD, C_DIV)
    rr(s, l, t, w, Cm(1.05), tc)
    txt(s, l+Cm(0.5), t+Cm(0.15), w-Cm(1.0), Cm(0.85), title, Pt(17), True, C_WHITE)
    mtxt(s, l+Cm(0.5), t+Cm(1.35), w-Cm(1.0), h-Cm(1.7), body_lines, bs, C_GRAY, sp=1.3)
    return s

def metric_box(s, l, t, w, h, value, label, color=C_BLUE):
    rr(s, l, t, w, h, C_CARD, C_DIV)
    txt(s, l, t+Cm(0.3), w, Cm(1.8), value, Pt(36), True, color, PP_ALIGN.CENTER)
    txt(s, l, t+h-Cm(1.5), w, Cm(1.2), label, Pt(13), False, C_GRAY, PP_ALIGN.CENTER)

def add_table(s, l, t, col_w, headers, rows, hc=C_DARK):
    ncols = len(headers); nrows = len(rows)+1; tw = sum(col_w)
    ts = s.shapes.add_table(nrows, ncols, Emu(l), Emu(t), Emu(tw), Emu(Cm(0.95)*nrows))
    tbl = ts.table
    for j, w in enumerate(col_w): tbl.columns[j].width = Emu(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = ""; c.fill.solid(); c.fill.fore_color.rgb = hc
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h; font(r, Pt(13), True, C_WHITE)
    for i, row in enumerate(rows):
        bg = C_CARD if i%2==0 else C_BG
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j); c.text = ""; c.fill.solid(); c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val); font(r, Pt(13), j==0, C_DARK)
    return ts

def arrow_r(s, l, t, w, h, color=C_ACC):
    shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(l), Emu(t), Emu(w), Emu(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()

def img_with_note(s, path, l, t, iw, ih, note_title, note_lines, note_x=None, note_w=Cm(15)):
    """Image on left + annotation on right (or vice versa)."""
    img(s, path, l, t, iw, ih)
    if note_x is None:
        note_x = l + iw + Cm(0.8)
    # Annotation card
    rr(s, note_x, t, note_w, ih, C_BG, C_DIV)
    mtxt(s, note_x+Cm(0.4), t+Cm(0.3), note_w-Cm(0.8), ih-Cm(0.6),
         [[note_title, True, C_BLUE, Pt(16)], ""] + note_lines,
         Pt(13), C_GRAY, sp=1.3)


# ══════════════════════════ BUILD ══════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    blank = prs.slide_layouts[6]
    TOTAL = 20

    # ── S1: COVER ────────────────────────────────────
    s = prs.slides.add_slide(blank)
    rr(s, Cm(0), Cm(0), SW, SH, C_DARK)
    oval(s, Cm(24), Cm(-3), Cm(16), Cm(16), RGBColor(0x24,0x38,0x62))
    oval(s, Cm(-5), Cm(7), Cm(12), Cm(12), RGBColor(0x24,0x38,0x62))
    rr(s, Cm(2.5), Cm(14.0), Cm(6), Cm(0.07), C_ACC)
    txt(s, Cm(2.5), Cm(3.0), Cm(29), Cm(3.0), "空调用电负荷预测", Pt(50), True, C_WHITE)
    txt(s, Cm(2.5), Cm(6.5), Cm(29), Cm(1.8), "基于多算法对比的时间序列预测实验", Pt(26), color=C_SUB)
    txt(s, Cm(2.5), Cm(9.5), Cm(29), Cm(1.2), "重庆大学  大数据与软件学院  ·  数据挖掘实验", Pt(18), color=C_LGRAY)
    txt(s, Cm(2.5), Cm(16.0), Cm(15), Cm(1.0), "汇报人：黄翰林    2026年6月", Pt(16), color=C_LGRAY)
    txt(s, Cm(2.5), Cm(17.2), Cm(20), Cm(0.8), "UCI Individual Household Electric Power Consumption", Pt(12), color=C_LGRAY)
    txt(s, Cm(30), Cm(17.8), Cm(3), Cm(0.5), "1", Pt(14), color=C_LGRAY, align=PP_ALIGN.RIGHT)

    # ── S2: TOC ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 2, "目  录", "CONTENTS")
    add_bottom(s)
    toc = [
        ("01", "实验背景与目标", "数据集概况、问题定义与研究目标"),
        ("02", "数据预处理与特征工程", "缺失值处理、归一化、时间特征、滞后与滚动统计"),
        ("03", "模型构建与训练", "Random Forest / XGBoost / LSTM / MLP-RF 四类算法详解"),
        ("04", "关键发现：数据泄漏排查", "从 R²=0.999 到物理公式泄漏的完整排查过程"),
        ("05", "实验结果与可视化分析", "性能对比、预测曲线、特征重要性、融合策略"),
        ("06", "总结与展望", "核心结论、不足分析与改进方向"),
    ]
    for i, (num, title, desc) in enumerate(toc):
        y = Cm(3.3) + Cm(2.5)*i
        oval(s, Cm(2.5), y+Cm(0.05), Cm(1.4), Cm(1.4), C_DARK)
        txt(s, Cm(2.5), y+Cm(0.2), Cm(1.4), Cm(1.1), num, Pt(22), True, C_WHITE, PP_ALIGN.CENTER)
        txt(s, Cm(4.5), y+Cm(0.05), Cm(25), Cm(1.1), title, Pt(24), True, C_DARK)
        txt(s, Cm(4.5), y+Cm(1.2), Cm(25), Cm(0.7), desc, Pt(14), color=C_GRAY)

    # ── S3: Background & Goals ──────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 3, "01  实验背景与目标", "Background & Objectives")
    add_bottom(s)

    card(s, Cm(1.5), Cm(3.4), Cm(15), Cm(6.5),
         "📋 实验背景", [
             "• 电力负荷预测是智能电网与能源管理的核心问题",
             "• 准确的短期负荷预测可优化发电调度、降低运营成本",
             "• 空调负荷在家庭用电中占比最高（约40%），季节性波动显著",
             "• 传统单一模型难以同时捕捉负荷的多尺度时序特征",
             "• 使用 UCI 家庭用电数据集（47个月，200万+条原始记录）",
         ], C_BLUE, Pt(14))

    card(s, Cm(17.5), Cm(3.4), Cm(15), Cm(6.5),
         "🎯 实验目标", [
             "① 实现完整的数据预处理与特征工程流水线",
             "② 构建并调优 4 种不同范式的预测模型",
             "③ 使用 5 项指标（MSE/MAE/RMSE/MAPE/R²）全面评估",
             "④ 通过可视化与模型诊断进行深度对比分析",
             "⑤ 探究数据泄漏、特征重要性、融合策略等关键问题",
         ], C_ACC, Pt(14))

    mtxt(s, Cm(1.5), Cm(10.5), Cm(31), Cm(7.0), [
        ("▎数据集概况", True, C_DARK, Pt(20)),
        "",
        "UCI Individual Household Electric Power Consumption — 法国巴黎某家庭 2006.12 ~ 2010.11 的分钟级能耗数据",
        "",
        ("• 7个变量：", True, C_DARK, Pt(15)),
        "  Global_active_power（目标变量·总有功功率）、Global_reactive_power、Voltage、Global_intensity、Sub_metering_1/2/3",
        ("• 数据规模：", True, C_DARK, Pt(15)),
        "  原始 2,075,259 行 → 逐小时重采样 34,168 行 → 去除NaN后 34,144 行（~3.9年逐时数据）",
        ("• 数据划分：", True, C_DARK, Pt(15)),
        "  训练 70% (2006.12-2009.09) | 验证 15% | 测试 15% (2010.04-2010.11) — 严格时序分割，防止未来信息泄漏",
    ], Pt(14), C_GRAY)

    # ── S4: Preprocessing Pipeline ───────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 4, "02  数据预处理流水线", "Data Preprocessing Pipeline")
    add_bottom(s)

    steps_text = [
        ("原始数据加载", "127MB .txt\nfloat32内存优化\n仅加载所需列"),
        ("缺失值处理", "1.25%行含NaN\n直接删除\n避免引入偏差"),
        ("逐小时重采样", "2,075,259→34,168行\n取每小时均值\n大幅压缩数据"),
        ("归一化", "Min-Max→[0,1]\n消除量纲差异\n加速梯度下降"),
        ("特征构建", "时间循环编码\n6个滞后特征\n12个滚动统计"),
        ("时序分割", "70/15/15\n严格时序切分\n防未来信息泄漏"),
    ]
    x0 = Cm(0.8)
    for i, (title, desc) in enumerate(steps_text):
        x = x0 + Cm(5.5)*i
        oval(s, x+Cm(1.5), Cm(3.1), Cm(1.1), Cm(1.1), C_ACC)
        txt(s, x+Cm(1.5), Cm(3.2), Cm(1.1), Cm(0.9), str(i+1), Pt(18), True, C_WHITE, PP_ALIGN.CENTER)
        txt(s, x+Cm(0.1), Cm(4.5), Cm(5.0), Cm(1.0), title, Pt(16), True, C_DARK, PP_ALIGN.CENTER)
        mtxt(s, x+Cm(0.1), Cm(5.5), Cm(5.0), Cm(3.0), desc.split("\n"), Pt(12), C_GRAY, PP_ALIGN.CENTER, sp=1.3)
        if i < 5: arrow_r(s, x+Cm(5.0), Cm(3.3), Cm(0.5), Cm(0.5), C_DIV)

    card(s, Cm(1.5), Cm(9.0), Cm(15), Cm(8.0),
         "⏱ 时间特征与循环编码", [
             "【原始时间特征】hour、dayofweek、month、is_weekend",
             "",
             "【sin/cos 循环编码 — 保留周期性】",
             "  hour_sin = sin(2π·hour/24)，hour_cos = cos(2π·hour/24)",
             "  → 23时与0时在圆上相邻，而非数值跳变",
             "  month_sin = sin(2π·month/12)，month_cos = cos(2π·month/12)",
             "  → 12月与1月在圆上相邻，保留年周期循环",
         ], C_BLUE, Pt(13))

    card(s, Cm(17.5), Cm(9.0), Cm(15), Cm(8.0),
         "📐 滞后特征与滚动统计", [
             "【滞后特征 — 6个，编码时序自相关】",
             "  lag1h/2h/3h/6h/12h/24h → 前一小时负荷是最直接的预测信息",
             "",
             "【滚动窗口统计 — 12个】",
             "  窗口 6h/12h/24h × 统计量 mean/std/min/max",
             "  roll_mean_24h = 过去一天的平均负荷水平",
             "  roll_std_6h = 近6小时的波动程度",
             "",
             "【最终：30维特征向量 × 34,144条样本】",
         ], C_ACC, Pt(13))

    # ── S5: EDA ──────────────────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 5, "02  探索性数据分析（EDA）", "Exploratory Data Analysis")
    add_bottom(s)

    # Two images with annotations
    img_with_note(s, os.path.join(FIGURES, "eda_correlation_heatmap.png"),
                  Cm(0.5), Cm(3.2), Cm(15.5), Cm(7.0),
                  "▼ 变量相关性热力图",
                  ["展示7个原始变量间的Pearson相关系数。",
                   "Global_active_power 与 Global_intensity 的相关系数≈1.0 —",
                   "由物理公式P=V×I决定，验证了排除该特征的合理性。",
                   "Sub_metering_3（空调）与目标相关性最高(~0.7)，是核心预测变量。"],
                  note_x=Cm(16.5), note_w=Cm(17.0))

    img_with_note(s, os.path.join(FIGURES, "eda_seasonal_decompose.png"),
                  Cm(0.5), Cm(10.8), Cm(15.5), Cm(6.8),
                  "▼ 季节性分解（周期7天）",
                  ["将逐日均值序列分解为趋势(Trend)、季节性(Seasonal)和残差(Residual)。",
                   "长期趋势在2008-2009年上升后回落；周季节性模式明显（周末用电>工作日）；",
                   "残差基本呈白噪声特征，说明周期7天的分解有效提取了主要规律。"],
                  note_x=Cm(16.5), note_w=Cm(17.0))

    # ── S6: Model Overview ───────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 6, "03  模型架构总览", "Model Architecture Overview")
    add_bottom(s)

    models = [
        ("🌲 随机森林\nRandom Forest", "Bagging 集成\n200棵树 · max_depth=25\n抗噪 · 不易过拟合\n训练: 1201s", C_BLUE),
        ("⚡ XGBoost", "Boosting 集成\n400棵树 · max_depth=6\n正则化 · 自动处理缺失\n训练: 106s", C_GREEN),
        ("🧠 LSTM", "3层循环神经网络\nhidden_dim=128\ndropout=0.3 · 梯度裁剪\n早停: patience=15", C_PURPLE),
        ("🔗 MLP-RF 融合", "MLP(256→128→64)\n+ RF 简单平均\nBatchNorm + Dropout\n训练: 23s (MLP部分)", C_ACC),
    ]
    for i, (name, desc, clr) in enumerate(models):
        x = Cm(1.0) + Cm(8.2)*i; w = Cm(7.8)
        rr(s, x, Cm(3.2), w, Cm(2.8), clr)
        txt(s, x+Cm(0.3), Cm(3.5), w-Cm(0.6), Cm(2.3), name, Pt(18), True, C_WHITE, PP_ALIGN.CENTER)
        rr(s, x, Cm(6.0), w, Cm(7.5), C_CARD, C_DIV)
        mtxt(s, x+Cm(0.5), Cm(6.3), w-Cm(1.0), Cm(6.8), desc.split('\n'), Pt(15), C_GRAY, sp=1.5)

    txt(s, Cm(2), Cm(14.5), Cm(30), Cm(3.0),
        "训练环境：PyTorch 2.12 + CUDA 13.0  |  NVIDIA RTX 4060 Laptop (8GB)  |  "
        "超参搜索：RandomizedSearchCV / 3折交叉验证  |  "
        "评估指标：MSE / MAE / RMSE / MAPE / R²",
        Pt(12), color=C_LGRAY, align=PP_ALIGN.CENTER)

    # ── S7: RF Detail ────────────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 7, "03  模型详解：随机森林", "Random Forest — Bagging Ensemble")
    add_bottom(s)
    card(s, Cm(1.5), Cm(3.4), Cm(15), Cm(13.5),
         "算法原理与实现", [
             "【Bagging 集成学习】并行构建多个决策树，取预测均值",
             "  两个随机性：① 样本Bootstrap采样 ② 分裂时随机特征子集",
             "  通过增加多样性来降低方差，天然抗过拟合",
             "",
             "【本次实验做了什么】",
             "  · 使用 sklearn RandomForestRegressor 构建模型",
             "  · RandomizedSearchCV：15轮 × 3折交叉验证，搜索5个关键超参数",
             "  · 在训练+验证集上用最优参数重新训练",
             "",
             "【核心优势】",
             "  ✓ 无需特征归一化（树模型不受量纲影响）",
             "  ✓ 可输出特征重要性，辅助理解数据",
             "  ✓ 对异常值不敏感，天然抗过拟合",
         ], C_BLUE, Pt(14))
    card(s, Cm(17.5), Cm(3.4), Cm(15), Cm(13.5),
         "最优超参数 & 关键指标", [
             ("【最优超参数】", True, C_ACC, Pt(15)),
             "  n_estimators=200  max_depth=25  min_samples_split=5",
             "  min_samples_leaf=4  max_features=0.5",
             "",
             ("【测试集性能】", True, C_ACC, Pt(15)),
             "  MSE=0.0013  MAE=0.0238  RMSE=0.0359",
             "  MAPE=25.09%  R²=0.8937",
             "",
             ("【训练效率】训练1201s（单线程），预测<1s", True, C_DARK, Pt(14)),
         ], C_ACC, Pt(14))

    # ── S8: XGBoost Detail ───────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 8, "03  模型详解：XGBoost", "XGBoost — Gradient Boosting Ensemble")
    add_bottom(s)
    card(s, Cm(1.5), Cm(3.4), Cm(15), Cm(13.5),
         "算法原理与实现", [
             "【Gradient Boosting】树串行构建，每棵新树拟合前一棵树的残差",
             "  目标函数 = 损失函数 + 正则项Ω（控制模型复杂度）",
             "  tree_method='hist'：基于直方图的近似分裂 → 大幅加速",
             "  内置缺失值自动学习 + 反向剪枝（避免局部最优）",
             "",
             "【本次实验做了什么】",
             "  · XGBRegressor(tree_method='hist') 加速训练",
             "  · RandomizedSearchCV：20轮 × 3折 CV，搜索7个超参数",
             "  · L1(reg_alpha)+L2(reg_lambda)正则联合控制过拟合",
             "  · subsample+colsample双重采样增加随机性",
         ], C_GREEN, Pt(14))
    card(s, Cm(17.5), Cm(3.4), Cm(15), Cm(13.5),
         "最优超参数 & 关键指标", [
             ("【最优超参数】", True, C_ACC, Pt(15)),
             "  n_estimators=400  max_depth=6  lr=0.05",
             "  subsample=0.8  colsample_bytree=0.7",
             "  reg_alpha=0  reg_lambda=1.5",
             "",
             ("【测试集性能】", True, C_ACC, Pt(15)),
             "  MSE=0.0013  MAE=0.0232  RMSE=0.0358",
             "  MAPE=23.79% ← 所有模型中最低！  R²=0.8947",
             "",
             ("⚡ 训练仅106s（RF的1/11），兼顾精度与效率，实际部署首选", True, C_RED, Pt(15)),
         ], C_ACC, Pt(14))

    # ── S9: LSTM ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 9, "03  模型详解：LSTM", "Long Short-Term Memory Network")
    add_bottom(s)

    card(s, Cm(1.5), Cm(3.4), Cm(15), Cm(6.5),
         "网络结构与训练配置", [
             "【3层 LSTM (PyTorch)】hidden_dim=128, num_layers=3, dropout=0.3",
             "  输入：(batch, 24h, 10维) → 过去24小时的10维多变量序列",
             "  回归头：Linear(128→64)+ReLU+Dropout+Linear(64→1)",
             "  参数量 ~200K，输出下一小时负荷预测值",
             "",
             "【训练策略】Adam(lr=1e-3) + ReduceLROnPlateau学习率衰减",
             "  梯度裁剪max_norm=1.0 + 早停patience=15 + Batch Size=64",
         ], C_PURPLE, Pt(14))

    card(s, Cm(1.5), Cm(10.8), Cm(15), Cm(6.3),
         "训练过程与性能", [
             "第48轮达最优验证MSE≈0.0067，第63轮触发早停",
             "训练loss持续下降但验证loss停滞 → 模型已达容量上限",
             "验证MSE远高于RF/XGBoost的~0.0013，预示最终性能差距",
             "",
             "测试集：MSE=0.0051  MAE=0.0489  RMSE=0.0715",
             "            MAPE=53.92%  R²=0.5793",
         ], C_DARK, Pt(14))

    # Image with annotation
    img_with_note(s, os.path.join(FIGURES, "lstm_training_curves.png"),
                  Cm(17.0), Cm(3.4), Cm(16), Cm(6.5),
                  "▼ LSTM训练曲线解读",
                  ["训练损失（蓝线）持续下降，但验证损失（橙线）",
                   "在第48轮后不再改善并在~0.007处停滞。",
                   "这说明模型已学到数据中能学到的所有模式，",
                   "继续训练只会过拟合。早停机制在第63轮",
                   "（连续15轮无改善后）终止训练。"],
                  note_x=Cm(17.0), note_w=Cm(16.0))

    # Bottom insight for LSTM slide
    rr(s, Cm(17.0), Cm(10.5), Cm(16), Cm(6.5), C_BG, C_DIV)
    mtxt(s, Cm(17.5), Cm(10.8), Cm(15.0), Cm(6.0), [
        ("▎为什么LSTM效果不理想？", True, C_RED, Pt(16)),
        "",
        "LSTM使用10维原始序列，而RF/XGBoost获得30维",
        "显式特征。LSTM需要从序列中隐式学习滞后和滚动",
        "模式，对模型容量和数据量要求更高。在~29K样本",
        "的规模下，树模型的结构先验优势更加明显。",
    ], Pt(14), C_GRAY, sp=1.3)

    # ── S10: MLP-RF Ensemble ────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 10, "03  模型详解：MLP-RF融合", "MLP-RF Ensemble — 神经网络 + 树模型")
    add_bottom(s)

    card(s, Cm(1.5), Cm(3.4), Cm(15), Cm(9.0),
         "MLP 网络结构", [
             "【3层全连接网络（PyTorch）】",
             "  Linear(30→256)+BatchNorm+ReLU+Dropout(0.3)",
             "  Linear(256→128)+BatchNorm+ReLU+Dropout(0.3)",
             "  Linear(128→64)+BatchNorm+ReLU+Dropout(0.3)",
             "  Linear(64→1) → 输出预测值",
             "",
             "【设计考量】BatchNorm加速收敛、Dropout(0.3)防过拟合",
             "  30维显式特征降低了模型隐式学习的要求",
             "  训练仅23s，性能接近树模型：R²=0.8904",
         ], C_PURPLE, Pt(14))

    card(s, Cm(1.5), Cm(13.2), Cm(15), Cm(3.8),
         "关键启示", [
             "特征工程做得足够好时，简单MLP也能接近复杂树模型的性能",
             "BatchNorm+Dropout是在表格数据上训练NN的关键组合",
         ], C_DARK, Pt(14))

    card(s, Cm(17.5), Cm(3.4), Cm(15), Cm(13.6),
         "两种融合策略对比", [
             ("【策略A：Ridge Stacking → R²=0.8905 ✗】", True, C_RED, Pt(16)),
             "  以RF和MLP预测值为特征，Ridge回归作元学习器",
             "  Ridge给MLP的权重≈-0.07 → 融合退化为只用RF",
             "  失败原因：RF与MLP预测相关系数ρ>0.95，高度共线",
             "  Stacking假设基模型提供互补视角，过于相似时无信息可整合",
             "",
             ("【策略B：简单算术平均 → R²=0.9007 ✓】", True, C_GREEN, Pt(16)),
             "  ensemble = (RF_pred + MLP_pred) / 2",
             "  强制各50%权重 → 不同偏差模式部分抵消",
             "  R²从0.8937/0.8904提升至0.9007，整体优于部分",
             "",
             ("💡 基模型预测高度相关时，简单平均可能比复杂元学习更有效", True, C_ACC, Pt(16)),
         ], C_ACC, Pt(14))

    # ── S11: Data Leakage ────────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 11, "04  关键发现：数据泄漏排查", "Critical Discovery — Data Leakage Detection & Fix")
    add_bottom(s)
    rr(s, Cm(1.5), Cm(3.4), Cm(31), Cm(2.8), RGBColor(0xFD,0xED,0xEC), C_RED)
    mtxt(s, Cm(2.5), Cm(3.6), Cm(29), Cm(2.4), [
        ("⚠️  实验初期，模型 R² 一度达到 0.999 — 表面近乎完美", True, C_RED, Pt(22)),
        ("   特征重要性分析揭示：Global_intensity（电流）占据 99.9% 的重要性 — 这是数据泄漏的警报", False, C_DARK, Pt(16)),
    ], Pt(16))

    card(s, Cm(1.5), Cm(7.0), Cm(15), Cm(10.0),
         "根因分析", [
             ("物理定律：P = V × I（功率 = 电压 × 电流）", True, C_DARK, Pt(17)),
             "",
             "法国电网中电压 V ≈ 240V（几乎恒定），因此：",
             "  Global_active_power ≈ 240 × Global_intensity",
             "",
             "将电流作为特征 → 模型只需学会「乘以240」即得答案",
             "这不是「预测」，而是「读取答案」— 典型的数据泄漏",
             "",
             ("这不是模型的错，是我们的疏忽：", True, C_RED, Pt(15)),
             "给了模型一个和答案线性相关的输入，然后夸它准。",
         ], C_RED, Pt(14))

    card(s, Cm(17.5), Cm(7.0), Cm(15), Cm(10.0),
         "修复措施 & 核心教训", [
             ("【修复措施】", True, C_GREEN, Pt(17)),
             "  排除 Global_intensity、Global_reactive_power",
             "  修复后 R² 从 0.999 → 0.89~0.90",
             "",
             ("【核心教训】", True, C_ACC, Pt(17)),
             "  ① R²=0.999 不是好消息，是数据泄漏的警报",
             "  ② 任一特征重要性>90%时，立刻审视其与目标的关系",
             "  ③ 理解数据背后的物理含义，比盲目追求高指标更重要",
             "  ④ ML第一法则：不要让模型看到答案",
         ], C_BLUE, Pt(14))

    # ── S12: Results Table ──────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 12, "05  实验结果：五模型性能对比", "Results — Model Performance Comparison")
    add_bottom(s)

    col_w = [Cm(5.5), Cm(4.5), Cm(4.5), Cm(4.5), Cm(5.0), Cm(5.0)]
    headers = ["模型 Model", "MSE ↓", "MAE ↓", "RMSE ↓", "MAPE(%) ↓", "R² ↑"]
    rows = [
        ["MLP-RF Avg", "0.0012", "0.0228", "0.0347", "24.60", "0.9007"],
        ["XGBoost",    "0.0013", "0.0232", "0.0358", "23.79", "0.8947"],
        ["Random Forest","0.0013","0.0238", "0.0359", "25.09", "0.8937"],
        ["MLP",         "0.0013", "0.0240", "0.0365", "26.15", "0.8904"],
        ["LSTM",        "0.0051", "0.0489", "0.0715", "53.92", "0.5793"],
    ]
    add_table(s, Cm(1.5), Cm(3.5), col_w, headers, rows)

    metrics_cards = [
        ("0.9007", "最佳R²\nMLP-RF Avg", C_GREEN),
        ("23.79%", "最佳MAPE\nXGBoost", C_BLUE),
        ("0.0012", "最低MSE\nMLP-RF Avg", C_GREEN),
        ("106s", "最快训练\nXGBoost", C_ACC),
    ]
    for i, (val, label, clr) in enumerate(metrics_cards):
        metric_box(s, Cm(1.5)+Cm(8.0)*i, Cm(9.5), Cm(7.2), Cm(3.2), val, label, clr)

    mtxt(s, Cm(1.5), Cm(13.5), Cm(31), Cm(4.5), [
        ("🏆 MLP-RF融合 R²=0.9007 — 简单平均超越所有单一模型，证明集成学习的价值", False, C_DARK, Pt(16)),
        ("⚡ XGBoost MAPE=23.79% — 兼顾精度与效率(106s)，实际部署场景的最优选择", False, C_DARK, Pt(16)),
        ("📉 LSTM MAPE=53.9%, R²=0.579 — 特征不对称+数据量不足，详后续分析", False, C_RED, Pt(16)),
        ("📊 树模型在前三占两席 — 丰富的显式特征工程是树模型领先的关键", False, C_GRAY, Pt(16)),
    ], Pt(14), C_GRAY, sp=1.3)

    # ── S13: Prediction Viz ──────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 13, "05  预测结果可视化", "Results — Prediction Visualization")
    add_bottom(s)

    img_with_note(s, os.path.join(FIGURES, "all_models_7days.png"),
                  Cm(0.3), Cm(3.0), Cm(17.0), Cm(7.0),
                  "▼ 测试集前7天（168小时）预测曲线",
                  ["灰色线为真实观测值，彩色线为各模型预测。",
                   "MLP-RF Avg(红线)和XGBoost(蓝线)最贴合真实曲线；",
                   "LSTM偏差明显偏大，尤其在早晚负荷高峰响应迟缓。",
                   "前4个模型均能较好地跟踪负荷的日周期波动。",
                   "一周内可观察到周末(后段)用电模式与工作日的差异。"],
                  note_x=Cm(17.8), note_w=Cm(15.5))

    img_with_note(s, os.path.join(FIGURES, "all_models_scatter.png"),
                  Cm(0.3), Cm(10.6), Cm(17.0), Cm(7.0),
                  "▼ 预测值 vs 真实值散点图",
                  ["每个点代表一个测试样本。横轴为真实值，纵轴为预测值。",
                   "越接近对角线y=x表示预测越准确。",
                   "树模型(RF/XGBoost)和MLP的点紧密围绕对角线；",
                   "LSTM在低值区和高值区均存在较大偏差；",
                   "MLP-RF融合的散点最为集中，与R²最高一致。"],
                  note_x=Cm(17.8), note_w=Cm(15.5))

    # ── S14: Feature Importance ──────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 14, "05  特征重要性分析", "Results — Feature Importance Analysis")
    add_bottom(s)

    mtxt(s, Cm(1.5), Cm(3.4), Cm(31), Cm(14.0), [
        ("RF 与 XGBoost 的特征重要性排序高度一致：", False, C_GRAY, Pt(16)),
        "",
        ("① Sub_metering_3（空调/热水器能耗）→ 重要性 ~38%  🔥", True, C_RED, Pt(20)),
        "    空调是家庭电力负荷的最强驱动因素 — 分项计量数据是预测的核心信号",
        "",
        ("② lag1h（前一小时负荷）→ 重要性 ~12-21%", True, C_BLUE, Pt(20)),
        "    电力负荷具有强时序自相关性 — 最近时刻的负荷是最直接的预测信息",
        "",
        ("③ Sub_metering_1 + Sub_metering_2（厨房+洗衣间）→ 合计约 20%", True, C_BLUE, Pt(20)),
        "    特定电器使用模式（做饭时间、洗衣时间）提供了额外预测信号",
        "",
        ("④ 时间循环特征（hour_sin/cos, month_cos）→ 合计约 5%", True, C_GRAY, Pt(18)),
        "    日周期和年周期特征为模型提供时间上下文信息",
        "",
        ("⑤ Voltage（电压）→ 2-3%", True, C_GRAY, Pt(18)),
        "    电网电压的微小波动也包含一定的预测信息",
        "",
        ("💡 结论：显式特征工程（滞后+滚动+循环编码）直接编码了时序依赖，是树模型领先的关键。", True, C_ACC, Pt(18)),
    ], Pt(14), C_GRAY, sp=1.1)

    # ── S15: Why LSTM failed ─────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 15, "05  深度分析：LSTM为何表现不佳？", "Analysis — Why LSTM Underperformed")
    add_bottom(s)
    reasons = [
        ("❶ 特征不对称（最核心原因）", C_RED, [
            "• LSTM：仅使用10维原始特征序列（含自回归项）",
            "• RF/XGBoost：获得30维显式工程特征（6滞后+12滚动+8时间−2泄漏）",
            "• LSTM需要从序列中「隐式学习」这些模式 → 对模型容量和数据量要求远超树模型",
            "• 这是我们实验设计的短板：为树模型做了充分特征工程，却未给LSTM同等待遇",
        ]),
        ("❷ 数据量不足", C_BLUE, [
            "• ~29K训练样本 对 3层 LSTM（~200K参数）偏少",
            "• 基于树的集成方法在小样本场景下天然具有结构上的先验优势",
            "• 深度学习通常需要更大规模数据才能发挥表示学习能力",
        ]),
        ("❸ 任务特性决定了模型选择", C_GREEN, [
            "• 对于已完成充分特征工程的表格化时序预测，梯度提升树（XGBoost）是更优选择",
            "• LSTM的优势场景：原始长序列、强非线性多变量交互、大规模数据端到端学习",
        ]),
    ]
    for i, (title, clr, bullets) in enumerate(reasons):
        card(s, Cm(1.5), Cm(3.5)+Cm(4.8)*i, Cm(31), Cm(4.3), title, bullets, clr, Pt(15))

    # ── S16: More Viz ────────────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 16, "05  更多可视化分析", "Additional Visualizations")
    add_bottom(s)

    # Top row: heatmap + residuals
    img_with_note(s, os.path.join(FIGURES, "eda_power_heatmap.png"),
                  Cm(0.3), Cm(3.0), Cm(15.5), Cm(6.5),
                  "▼ 功率变化热力图",
                  ["以日期为行、小时为列，颜色深浅表示功率大小。",
                   "白天10:00-18:00功率显著高于夜间；",
                   "午间和傍晚出现两个用电高峰（做饭与回家后）。",
                   "周末用电模式与工作日明显不同。"],
                  note_x=Cm(16.5), note_w=Cm(17.0))

    img_with_note(s, os.path.join(FIGURES, "rf_xgb_residuals.png"),
                  Cm(0.3), Cm(10.2), Cm(15.5), Cm(6.5),
                  "▼ RF与XGBoost残差分布",
                  ["残差=真实值−预测值，理想情况应呈正态分布、以零为中心。",
                   "两个模型的残差均近似正态、均值≈0，",
                   "说明不存在系统性偏差（不会系统性地高估或低估）。",
                   "XGBoost残差分布略窄于RF，与MAPE更低一致。"],
                  note_x=Cm(16.5), note_w=Cm(17.0))

    # ── S17: Stacking vs Average ────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 17, "05  融合策略深度对比", "Ensemble Strategy — Stacking vs. Simple Average")
    add_bottom(s)

    card(s, Cm(1.5), Cm(3.5), Cm(15), Cm(13.0),
         "Ridge Stacking（元学习融合）", [
             ("R² = 0.8905 — 不如单独使用 RF (0.8937)", True, C_RED, Pt(18)),
             "",
             "【为什么失败？】",
             "  · Ridge元学习器给MLP的权重≈-0.07（≈零）",
             "  · 融合退化为几乎只使用RF，失去集成意义",
             "  · RF和MLP在测试集上的预测相关系数ρ>0.95",
             "  · 两个模型输出高度共线 → 元学习器无信息可整合",
             "",
             "【本质】Stacking假设基模型提供「互补视角」",
             "  基模型过于相似时，Stacking反而有害",
         ], C_RED, Pt(14))

    card(s, Cm(17.5), Cm(3.5), Cm(15), Cm(13.0),
         "简单算术平均 ✨", [
             ("R² = 0.9007 — 超越所有单一模型！", True, C_GREEN, Pt(18)),
             "",
             "ensemble = (RF_pred + MLP_pred) / 2",
             "",
             "【为什么更有效？】",
             "  · 强制给两个模型各50%权重，不做偏好选择",
             "  · RF和MLP偏差模式不完全重合 → 50/50平均使误差部分抵消",
             "  · 实现了真正的「集成学习」— 整体优于部分",
             "",
             ("💡 核心启示：", True, C_ACC, Pt(18)),
             "  当基学习器输出高度相关时，复杂元学习可能事与愿违。",
             "  简单平均因强制均分权重而获得更好的泛化效果。",
             "  「简单」有时比「复杂」更有效。",
         ], C_GREEN, Pt(14))

    txt(s, Cm(15.8), Cm(8.5), Cm(2.5), Cm(2.0), "vs", Pt(40), True, C_LGRAY, PP_ALIGN.CENTER)

    # ── S18: Summary ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 18, "06  实验总结", "Summary & Key Conclusions")
    add_bottom(s)

    conclusions = [
        ("1", "XGBoost — 实际部署首选", "MAPE=23.8%, 训练仅106s, 精度与效率的最佳平衡点", C_GREEN),
        ("2", "MLP-RF融合 — 最高R² (0.9007)", "简单平均超越所有单一模型, 证明集成学习的核心价值", C_BLUE),
        ("3", "数据泄漏 — 最有价值的教训", "R²=0.999 不是庆祝而是警报, 物理理解 > 高指标", C_RED),
        ("4", "特征工程 — 性能上限的决定因素", "Sub_metering_3(~38%) + lag1h(~21%) 是最强预测信号", C_ACC),
        ("5", "LSTM — 深度学习并非万能", "特征不对称+数据不足→MAPE=54%, 工具需匹配场景", C_PURPLE),
        ("6", "简单平均 > Stacking", "基学习器高度相关时, 简单集成策略可能更有效", C_GRAY),
    ]
    for i, (num, title, desc, clr) in enumerate(conclusions):
        y = Cm(3.3) + Cm(2.5)*i
        oval(s, Cm(2.0), y+Cm(0.15), Cm(1.1), Cm(1.1), clr)
        txt(s, Cm(2.0), y+Cm(0.25), Cm(1.1), Cm(0.9), num, Pt(18), True, C_WHITE, PP_ALIGN.CENTER)
        txt(s, Cm(3.8), y+Cm(0.05), Cm(28), Cm(1.0), title, Pt(22), True, clr)
        txt(s, Cm(3.8), y+Cm(1.1), Cm(28), Cm(0.8), desc, Pt(15), color=C_GRAY)

    # ── S19: Limitations & Future ────────────────────
    s = prs.slides.add_slide(blank)
    page_bar(s, 19, "06  不足与展望", "Limitations & Future Work")
    add_bottom(s)

    card(s, Cm(1.5), Cm(3.4), Cm(15), Cm(13.5),
         "当前不足", [
             "• LSTM特征不对称：未平等对待深度模型",
             "• 样本量有限：~29K条对深度模型偏少",
             "• 超参搜索粗糙：仅15-20轮RandomizedSearch",
             "• 未引入外部特征：气象数据（温度/湿度）",
             "  → 温度是空调负荷最直接的驱动因素",
             "• 未尝试Transformer、Informer等新架构",
             "• 仅做单步预测（1h ahead），非多步预测",
             "• 未量化预测不确定性（仅点估计）",
         ], C_RED, Pt(14))

    card(s, Cm(17.5), Cm(3.4), Cm(15), Cm(13.5),
         "改进方向", [
             "① 公平实验：为所有模型提供相同特征集",
             "",
             "② 引入气象数据：温度、湿度是空调负荷核心驱动",
             "",
             "③ 尝试新架构：Transformer/Informer/TimesFM",
             "",
             "④ 系统超参优化：Optuna/Bayesian Optimization",
             "",
             "⑤ 多步预测：从1步→24步，真正日负荷预测",
             "",
             "⑥ 不确定性量化：预测区间而非仅点估计",
             "",
             "⑦ 更大规模数据：使用多年多用户数据",
         ], C_GREEN, Pt(14))

    # ── S20: END ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    rr(s, Cm(0), Cm(0), SW, SH, C_DARK)
    oval(s, Cm(24), Cm(-3), Cm(16), Cm(16), RGBColor(0x24,0x38,0x62))
    oval(s, Cm(-5), Cm(8), Cm(10), Cm(10), RGBColor(0x24,0x38,0x62))
    rr(s, Cm(2.5), Cm(13.0), Cm(6), Cm(0.07), C_ACC)
    txt(s, Cm(2.5), Cm(4.0), Cm(29), Cm(3.5), "感谢聆听！", Pt(56), True, C_WHITE)
    txt(s, Cm(2.5), Cm(8.0), Cm(29), Cm(1.8), "Thank You", Pt(30), color=C_SUB)
    txt(s, Cm(2.5), Cm(10.8), Cm(29), Cm(1.2), "空调用电负荷预测 — 基于多算法对比的时间序列预测实验", Pt(18), color=C_LGRAY)
    txt(s, Cm(2.5), Cm(14.5), Cm(29), Cm(1.0), "汇报人：黄翰林  |  重庆大学 大数据与软件学院  |  数据挖掘实验", Pt(16), color=C_LGRAY)
    txt(s, Cm(2.5), Cm(16.0), Cm(29), Cm(1.2), "欢迎提问与交流", Pt(24), True, C_ACC)

    # ══════════════════════════════════════════════════
    prs.save(OUTPUT)
    print(f"✅  PPT saved: {OUTPUT}")
    print(f"   Size: {os.path.getsize(OUTPUT)/1024:.1f} KB, Slides: {TOTAL}")

if __name__ == "__main__":
    build()
